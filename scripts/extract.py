#!/usr/bin/env python3
"""finals-prepper Phase 1: 提取各章文件的文字，产出逐章 text bundle。

用法: python extract.py <资料目录>
读取 .final_prep/manifest.json，对每章文件按类型提取文字；
PPT 内嵌图片、独立图片文件记录路径供 AI 直接读图（不做 OCR）。
产物写 .final_prep/chapters/<章id>.json，PPT 内嵌图存 .final_prep/media/。
"""
import argparse
import json
import os

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}


def _save_image(blob, ext, media_dir, chapter_id, idx, kind):
    ext = (ext or "png").lower().lstrip(".")
    fn = f"{chapter_id}_{kind}{idx}.{ext}"
    fp = os.path.join(media_dir, fn)
    with open(fp, "wb") as f:
        f.write(blob)
    return fp


def extract_pptx(path, media_dir, chapter_id):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return None, "python-pptx 未安装", []
    prs = Presentation(path)
    parts = []
    img_paths = []
    pic_idx = 0
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    pic_idx += 1
                    img_paths.append(
                        _save_image(image.blob, image.ext, media_dir, chapter_id, pic_idx, "pic"))
                except Exception:
                    pass
                continue
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    texts.append(" | ".join(cells))
        if slide.has_notes_slide:
            nt = slide.notes_slide.notes_text_frame.text.strip()
            if nt:
                texts.append("[备注] " + nt)
        if texts:
            parts.append(f"[第{i}页]\n" + "\n".join(texts))
    return "\n\n".join(parts), None, img_paths


def extract_docx(path):
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        from docx.oxml.ns import qn
    except ImportError:
        return None, "python-docx 未安装"
    d = Document(path)
    parts = []
    for child in d.element.body.iterchildren():
        if child.tag == qn("w:p"):
            t = Paragraph(child, d).text.strip()
            if t:
                parts.append(t)
        elif child.tag == qn("w:tbl"):
            tbl = Table(child, d)
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
    return "\n".join(parts), None


def extract_pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return None, "pypdf 未安装"
    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            t = (page.extract_text() or "").strip()
        except Exception:
            t = ""
        if t:
            parts.append(f"[第{i}页]\n{t}")
    return "\n\n".join(parts), None


def main():
    ap = argparse.ArgumentParser(description="提取资料文字，产出逐章 bundle")
    ap.add_argument("root", help="资料目录")
    args = ap.parse_args()
    root = os.path.abspath(args.root)

    mp = os.path.join(root, ".final_prep", "manifest.json")
    with open(mp, encoding="utf-8") as f:
        manifest = json.load(f)

    outdir = os.path.join(root, ".final_prep", "chapters")
    media_dir = os.path.join(root, ".final_prep", "media")
    os.makedirs(outdir, exist_ok=True)
    os.makedirs(media_dir, exist_ok=True)

    for ch in manifest["chapters"]:
        texts = []
        images = []
        warnings = []
        for fp in ch["files"]:
            ext = os.path.splitext(fp)[1].lower()
            if ext in IMAGE_EXTS:
                images.append(fp)
                continue
            if ext in (".pptx", ".ppt"):
                t, err, imgs = extract_pptx(fp, media_dir, ch["id"])
                if imgs:
                    images.extend(imgs)
            elif ext in (".docx", ".doc"):
                t, err = extract_docx(fp)
            elif ext == ".pdf":
                t, err = extract_pdf(fp)
            else:
                continue
            if err:
                warnings.append(f"{os.path.basename(fp)}: {err}")
                continue
            if t:
                texts.append(f"===== 文件: {os.path.basename(fp)} =====\n{t}")

        bundle = {
            "chapter_id": ch["id"],
            "label": ch["label"],
            "text": "\n\n".join(texts),
            "images": images,
            "warnings": warnings,
        }
        out = os.path.join(outdir, f"{ch['id']}.json")
        with open(out, "w", encoding="utf-8") as f:
            json.dump(bundle, f, ensure_ascii=False, indent=2)
        print(f"{ch['id']} {ch['label']}: 文字 {len(texts)} 段, 图片 {len(images)} 张"
              + (f", 警告 {len(warnings)} 条" if warnings else ""))

    print(f"bundle -> {outdir}")


if __name__ == "__main__":
    main()
